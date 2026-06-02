"""Nới rộng ô text bị LibreOffice cấp quá hẹp (gây wrap tiêu đề/bullet).

Marp xuất PPTX editable qua LibreOffice → mỗi tiêu đề/bullet thành 1 ô text
riêng, rộng vừa khít chữ. Một số ô bị tính HỤT vài pixel → chữ tràn → xuống
dòng (wrap) đè lên dòng/bullet kế. Script này nới các ô đứng-một-mình ra sát
lề phải để chữ luôn đủ chỗ; BỎ QUA ô có hàng xóm bên phải cùng dòng (mảnh
inline-code, ô bảng cột trái) để không đè nhau.

Chạy SAU marp + add_links.py:
    .venv\\Scripts\\python.exe slides\\fix_boxes.py
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX = Path(__file__).resolve().parent / "SLIDES.pptx"
RIGHT_EDGE = 11750000   # lề phải mục tiêu (slide rộng 12,192,000 EMU)
GAP = 30000             # chừa khe trước hàng xóm bên phải


def is_chrome(t):
    t = t.strip()
    if t == "LOCAL LLM · RAG":
        return True
    if t.startswith("©"):
        return True
    if re.match(r"^\d+\s*/\s*\d+$", t):
        return True
    return False


def same_line(a, b):
    top = max(a.top, b.top)
    bot = min(a.top + a.height, b.top + b.height)
    return (bot - top) > 0.5 * min(a.height, b.height)


prs = Presentation(str(PPTX))
n = 0
for slide in prs.slides:
    boxes = [sh for sh in slide.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()
             and not getattr(sh, "has_table", False)
             and sh.left is not None and sh.width is not None]
    for b in boxes:
        if is_chrome(b.text_frame.text):
            continue
        neighbors = [o for o in boxes if o is not b
                     and o.left > b.left + 50000
                     and same_line(b, o)]
        if neighbors:
            limit = min(o.left for o in neighbors) - GAP
        else:
            limit = RIGHT_EDGE
        new_w = limit - b.left
        if new_w > b.width:
            b.width = Emu(int(new_w))
            b.text_frame.word_wrap = True
            n += 1

prs.save(str(PPTX))
print(f"[OK] noi rong {n} o text trong {PPTX.name}")
