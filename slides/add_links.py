"""Gắn lại hyperlink vào SLIDES.pptx.

Marp xuất PPTX editable qua LibreOffice → LibreOffice làm RỚT hyperlink
(link tải / link localhost chỉ còn là chữ, không bấm được).
Chạy script này NGAY SAU lệnh marp để gắn lại link bấm được:

    .venv\\Scripts\\python.exe slides\\add_links.py
"""
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).resolve().parent / "SLIDES.pptx"

# Chữ hiển thị trên slide  ->  URL bấm vào
LINKS = {
    "ollama.com/library": "https://ollama.com/library",
    "ollama.com/download": "https://ollama.com/download",
    "docker.com/products/docker-desktop": "https://www.docker.com/products/docker-desktop/",
    "brew.sh": "https://brew.sh",
    "localhost:3000": "http://localhost:3000",
    "localhost:8080": "http://localhost:8080",
}

prs = Presentation(str(PPTX))
n = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                url = LINKS.get(run.text.strip())
                if url:
                    run.hyperlink.address = url
                    n += 1
prs.save(str(PPTX))
print(f"[OK] da gan {n} hyperlink vao {PPTX.name}")
